import streamlit as st
import pandas as pd
import re
import matplotlib.pyplot as plt
import seaborn as sns
from sklearn.feature_extraction.text import TfidfVectorizer
import numpy as np
from sentence_transformers import SentenceTransformer
from sklearn.cluster import KMeans
from nltk.tokenize import sent_tokenize
from sklearn.metrics.pairwise import cosine_similarity
from pptx import Presentation
from io import BytesIO
import requests
import nltk

nltk.download('punkt')

# Skill keywords list (expand as you like)
SKILL_KEYWORDS = [
    "python", "sql", "excel", "power bi", "tableau", "aws", "azure", "gcp",
    "r", "java", "c++", "javascript", "html", "css",
    "machine learning", "deep learning", "nlp", "data analysis",
    "spark", "hadoop", "kubernetes", "docker", "git", "linux"
]

from dotenv import load_dotenv
import os

load_dotenv()  # reads .env file

APP_ID = os.getenv("ADZUNA_APP_ID")
APP_KEY = os.getenv("ADZUNA_APP_KEY")

def extract_skills(text):
    found = []
    text_lower = text.lower()
    for skill in SKILL_KEYWORDS:
        if re.search(rf"\b{re.escape(skill)}\b", text_lower):
            found.append(skill.title())
    return list(set(found))

def fetch_adzuna_jobs(keyword, country, pages, results_per_page=50):
    all_jobs = []
    for page in range(1, pages + 1):
        endpoint = f"https://api.adzuna.com/v1/api/jobs/{country}/search/{page}"
        params = {
            "app_id": APP_ID,
            "app_key": APP_KEY,
            "results_per_page": results_per_page,
            "what": keyword,
            "content-type": "application/json"
        }
        resp = requests.get(endpoint, params=params)
        if resp.status_code == 200:
            data = resp.json()
            all_jobs.extend(data.get("results", []))
        else:
            st.error(f"API Error (page {page}): {resp.status_code}")
            break

    cleaned = []
    for job in all_jobs:
        desc = job.get("description", "").strip()
        cleaned.append({
            "job_title": job.get("title", "").strip(),
            "company": job.get("company", {}).get("display_name", "").strip(),
            "location": job.get("location", {}).get("display_name", "").strip(),
            "description": desc,
            "skills": extract_skills(desc)
        })
    df = pd.DataFrame(cleaned)
    df.drop_duplicates(subset=["job_title", "company", "location", "description"], inplace=True)
    df.reset_index(drop=True, inplace=True)
    return df

def top_skills_frequency(df, min_count=5):
    counts = {}
    for skills in df['skills']:
        for skill in skills:
            counts[skill] = counts.get(skill, 0) + 1
    freq_df = pd.DataFrame(counts.items(), columns=["Skill", "Count"])
    freq_df = freq_df[freq_df["Count"] >= min_count].sort_values(by="Count", ascending=False)
    freq_df.reset_index(drop=True, inplace=True)
    return freq_df

def plot_bar_chart(freq_df, top_n=15):
    fig, ax = plt.subplots(figsize=(10,6))
    sns.barplot(data=freq_df.head(top_n), y='Skill', x='Count', palette='magma', ax=ax)
    ax.set_title(f"Top {top_n} Skills")
    st.pyplot(fig)

def plot_heatmap(df, skills_list):
    locations = df['location'].unique()
    heat_data = pd.DataFrame(0, index=skills_list, columns=locations)
    for _, row in df.iterrows():
        loc = row['location']
        for skill in row['skills']:
            if skill in skills_list:
                heat_data.loc[skill, loc] += 1
    fig, ax = plt.subplots(figsize=(14,8))
    sns.heatmap(heat_data, annot=True, fmt='d', cmap='coolwarm', ax=ax)
    ax.set_title("Skill Demand Heatmap by Location")
    st.pyplot(fig)

def get_top_tfidf_keywords(docs, top_n=20):
    vectorizer = TfidfVectorizer(stop_words='english', max_features=5000)
    X = vectorizer.fit_transform(docs)
    indices = np.argsort(np.asarray(X.sum(axis=0)).ravel())[::-1]
    features = np.array(vectorizer.get_feature_names_out())
    top_features = features[indices][:top_n]
    top_scores = np.asarray(X.sum(axis=0)).ravel()[indices][:top_n]
    return list(zip(top_features, top_scores))

def bert_cluster_docs(docs, n_clusters=4):
    model = SentenceTransformer('all-MiniLM-L6-v2')
    embeddings = model.encode(docs, show_progress_bar=False)
    km = KMeans(n_clusters=n_clusters, random_state=42)
    clusters = km.fit_predict(embeddings)
    return clusters

def extractive_summary_tfidf(text, n_sentences=3):
    sentences = sent_tokenize(text)
    if len(sentences) <= n_sentences:
        return text
    vectorizer = TfidfVectorizer(stop_words='english')
    sentence_vectors = vectorizer.fit_transform(sentences)
    sim_matrix = cosine_similarity(sentence_vectors)
    scores = sim_matrix.sum(axis=1)
    ranked_sentences = [sentences[i] for i in np.argsort(scores)[::-1][:n_sentences]]
    summary = ' '.join(ranked_sentences)
    return summary

def create_pptx_report(freq_df):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "LinkedIn Job Market Analysis"
    slide.placeholders[1].text = "Generated with Python and Streamlit\n\n"

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Top Skills"
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame

    for idx, row in freq_df.head(10).iterrows():
        p = tf.add_paragraph()
        p.text = f"{row['Skill']}: {row['Count']} job postings"
        p.level = 0

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio

def main():
    st.title("Enhanced LinkedIn Job Market Analysis")

    keyword = st.text_input("Job Keyword", value="data scientist")
    country = st.text_input("Country Code (e.g. us, gb, ca)", value="us")
    pages = st.slider("Number of Pages to Fetch (50 jobs per page)", 1, 5, 2)

    if st.button("Fetch and Analyze Jobs"):
        with st.spinner("Fetching data..."):
            df = fetch_adzuna_jobs(keyword, country, pages)

        if df.empty:
            st.warning("No jobs found for this query.")
            return

        st.success(f"Fetched {len(df)} job postings")

        freq_df = top_skills_frequency(df)
        st.subheader("Top Skills Frequency")
        st.dataframe(freq_df)

        plot_bar_chart(freq_df)

        selected_skills = st.multiselect("Select skills for heatmap", freq_df['Skill'].tolist(), default=freq_df['Skill'].head(5).tolist())
        if selected_skills:
            plot_heatmap(df, selected_skills)

        st.subheader("Top TF-IDF Keywords in Job Descriptions")
        keywords = get_top_tfidf_keywords(df['description'].tolist())
        st.write(pd.DataFrame(keywords, columns=['Keyword', 'Score']))

        st.subheader("Job Description Clustering (BERT + KMeans)")
        clusters = bert_cluster_docs(df['description'].tolist())
        df['cluster'] = clusters
        st.dataframe(df[['job_title', 'company', 'location', 'cluster']].head(10))

        st.subheader("Sample Extractive Summary")
        sample_idx = st.slider("Select Job Description Index to Summarize", 0, len(df)-1, 0)
        summary = extractive_summary_tfidf(df.loc[sample_idx, 'description'])
        st.write(summary)

        st.subheader("Export Slide Deck")
        pptx_data = create_pptx_report(freq_df)
        st.download_button(
            label="Download PPTX Report",
            data=pptx_data,
            file_name="job_market_analysis.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

if __name__ == "__main__":
    main()
# Run the Streamlit app
# To run this app, save it as app.py and use the command: streamlit run app.py
# Ensure you have the required libraries installed: streamlit, pandas, matplotlib, seaborn, sklearn, sentence-transformers, nltk, pptx
# You can install them using pip
# pip install streamlit pandas matplotlib seaborn scikit-learn sentence-transformers nltk python-pptx
# Make sure to have the nltk punkt tokenizer downloaded as well
# nltk.download('punkt')
# This app will fetch job postings, analyze skills, visualize data, and generate a PowerPoint report
# It uses the Adzuna API to fetch job data, so ensure you have a valid APP_ID and APP_KEY
# You can replace the APP_ID and APP_KEY with your own  credentials from Adzuna
# The app provides interactive visualizations and allows users to explore job market trends
# It also includes NLP features like TF-IDF keyword extraction and BERT clustering for job descriptions
# The app is designed to be user-friendly and provides a comprehensive analysis of the job market
# You can customize the job keyword and country code to fetch relevant job postings
# The app is built using Streamlit, which allows for easy deployment and sharing of the analysis
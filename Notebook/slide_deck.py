from pptx import Presentation

def create_slide_deck(skills_freq_df, save_path="job_market_analysis.pptx"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "LinkedIn Job Market Analysis"

    # Check if placeholder 1 exists and set text safely
    try:
        slide.placeholders[1].text = "Generated with Python\n\n"
    except IndexError:
        # fallback: add a textbox if placeholder doesn't exist
        from pptx.util import Inches, Pt
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(1.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = "Generated with Python\n\n"


    slide = prs.slides.add_slide(prs.slide_layouts[1])  # use Title and Content layout
    slide.shapes.title.text = "Top Skills in Job Market"
    body_shape = slide.shapes.placeholders[1]  # content placeholder
    tf = body_shape.text_frame

    for idx, row in skills_freq_df.head(10).iterrows():
        p = tf.add_paragraph()
        p.text = f"{row['Skill']}: {row['Count']} postings"
        p.level = 0

    prs.save(save_path)
    print(f"Slide deck saved to {save_path}")
